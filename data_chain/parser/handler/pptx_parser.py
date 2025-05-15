import os
from pptx import Presentation
from pptx.table import Table
import uuid
from data_chain.entities.enum import DocParseRelutTopology, ChunkParseTopology, ChunkType
from data_chain.parser.parse_result import ParseNode, ParseResult
from data_chain.parser.handler.base_parser import BaseParser
from data_chain.logger.logger import logger as logging


class PptxParser(BaseParser):
    name = 'pptx'

    @staticmethod
    async def extract_table_to_array(table: Table) -> list[list[str]]:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = ''.join([p.text for p in cell.text_frame.paragraphs])
                row_data.append(cell_text)
            table_data.append(row_data)
        return table_data

    @staticmethod
    async def ppt_to_parse_nodes(pptx) -> list[ParseNode]:
        nodes = []

        for slide_num, slide in enumerate(pptx.slides, start=1):
            for shape in slide.shapes:
                # 提取文字
                if shape.has_text_frame:
                    text = ""
                    try:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += run.text
                    except Exception as e:
                        err = "文字提取失败"
                        logging.exception("[PptxParser] %s", err)
                    if text.strip():
                        nodes.append(
                            ParseNode(
                                id=uuid.uuid4(),

                                lv=0,
                                parse_topology_type=ChunkParseTopology.GERNERAL,
                                content=text,
                                type=ChunkType.TEXT,
                                link_nodes=[]
                            )
                        )
                # 提取表格
                elif shape.has_table:
                    table = shape.table
                    table_array = await PptxParser.extract_table_to_array(table)
                    for row in table_array:
                        node = ParseNode(
                            id=uuid.uuid4(),

                            lv=0,
                            parse_topology_type=ChunkParseTopology.GERNERAL,
                            content=row,
                            type=ChunkType.TABLE,
                            link_nodes=[]
                        )
                        nodes.append(node)
                # 提取图片
                elif shape.shape_type == 13:  # 13 表示图片类型
                    try:
                        image = shape.image
                        blob = image.blob
                    except Exception as e:
                        err = "图片提取失败"
                        logging.exception("[PptxParser] %s", err)
                        continue
                    nodes.append(
                        ParseNode(
                            id=uuid.uuid4(),

                            lv=0,
                            parse_topology_type=ChunkParseTopology.GERNERAL,
                            content=blob,
                            type=ChunkType.IMAGE,
                            link_nodes=[]
                        )
                    )

        return nodes

    @staticmethod
    async def parser(file_path):
        try:
            pptx = Presentation(file_path)
        except Exception as e:
            err = "PPTX文件解析失败"
            logging.exception("[PptxParser] %s", err)
        nodes = await PptxParser.ppt_to_parse_nodes(pptx)
        PptxParser.image_related_node_in_link_nodes(nodes)
        parse_result = ParseResult(
            parse_topology_type=DocParseRelutTopology.LIST,
            nodes=nodes
        )
        return parse_result
